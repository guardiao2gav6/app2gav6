from pptx import Presentation
from pptx.util import Inches
import io
import base64
import gerador_de_graficos_tabelas
from dados_gsheets import Dados
import streamlit as st

# Carregando dados da página
dados = Dados()
detalhes_tripulantes_df = dados.generate_detalhes_tripulantes_df()
meta_pilotos_df = dados.generate_meta_pilotos_df()
dados_pessoais_df = dados.get_dados_pessoais()
descidas_df = dados.get_descidas()
aeronaves_df = dados.get_aeronaves()
esforco_aereo_df = dados.get_esforco_aereo()

# Função para gerar o gráfico com Altair e salvar como PNG em memória
def generate_chart_image(chart):
    # Exportar o gráfico para PNG usando a API Vega
    png_output = io.BytesIO()
    chart.properties(width=850,
                     height=200).configure_legend(labelFontSize=10,
                                                  titleFontSize=14).configure_axis(domain=False,
                                                                                   ticks=False).configure_view(
        strokeWidth=0
    ).save(
        png_output, format='png', ppi=96
    )
    png_output.seek(0)

    return png_output


def generate_esforco_aereo_chart_image(chart):
    # Exportar o gráfico para PNG usando a API Vega
    png_output = io.BytesIO()
    chart.properties(width=850,
                     height=250).configure_legend(labelFontSize=10,
                                                  titleFontSize=14).configure_axis(domain=False,
                                                                                   ticks=False).configure_view(
        strokeWidth=0
    ).save(
        png_output, format='png', ppi=96
    )
    png_output.seek(0)

    return png_output


def generate_charts_dict():

    pau_de_sebo_chart = gerador_de_graficos_tabelas.gerar_grafico_pau_de_sebo_impressao(
        detalhes_tripulantes_df=detalhes_tripulantes_df,
        meta_pilotos_df=meta_pilotos_df,
        dados_pessoais_df=dados_pessoais_df)[0]

    adaptacao_chart = gerador_de_graficos_tabelas.gerar_grafico_adaptacao_impressao(
        detalhes_tripulantes_df=detalhes_tripulantes_df,
        dados_pessoais_df=dados_pessoais_df
    )[0]

    funcoes_agrupadas = {'Mecânicos': ['AC', 'MC', 'IC'],
                         'Chefe Controlador': ['AB-R', 'CC-R', 'IB-R'],
                         'COTAT': ['AJ', 'CT', 'IJ'],
                         'COAM': ['AD-R', 'CO-R', 'ID-R'],
                         'O3': ['A3', 'O3', 'I3'],
                         'O1': ['A1', 'O1', 'I1'],
                         'MA-R': ['AA-R', 'MA-R', 'IA-R'],
                         'MA-E': ['AA-E', 'MA-E', 'IA-E'],
                         'Oficiais': ['AB-R', 'CC-R', 'IB-R', 'AJ', 'CT', 'IJ', '1P', '2P', 'IN', 'AL']}

    lista_funcoes_alunos = [i[0] for i in funcoes_agrupadas.values()]

    return {
        "Pau de Sebo - PIL": generate_chart_image(pau_de_sebo_chart),
        "Adaptação - PIL": generate_chart_image(adaptacao_chart),
        'Pau de Sebo - OFICIAIS': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='Oficiais',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - CC': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='Chefe Controlador',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - COTAT': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='COTAT',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - MC': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='Mecânicos',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - COAM': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='COAM',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - O1': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='O1',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - O3': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='O3',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - MA-E': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='MA-E',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
            dados_pessoais_df=dados_pessoais_df)[0]),
        'Pau de Sebo - MA-R': generate_chart_image(gerador_de_graficos_tabelas.gerar_grafico_demais_funcoes_impressao(
            funcao='MA-R',
            funcoes_agrupadas=funcoes_agrupadas,
            lista_funcoes_alunos=lista_funcoes_alunos,
            detalhes_tripulantes_df=detalhes_tripulantes_df,
                dados_pessoais_df=dados_pessoais_df)[0]),
        'Esforço Aéreo - E99 COMPREP': generate_esforco_aereo_chart_image(
            gerador_de_graficos_tabelas.gerar_grafico_esforco_aereo_impressao(
                aeronave='E-99',
                grupo='COMPREP',
                esforco_aereo_df=esforco_aereo_df)),
        'Esforço Aéreo - E99 COMAE': generate_esforco_aereo_chart_image(
            gerador_de_graficos_tabelas.gerar_grafico_esforco_aereo_impressao(
                aeronave='E-99',
                grupo='COMAE',
                esforco_aereo_df=esforco_aereo_df)),
        'Esforço Aéreo - E99 DCTA': generate_esforco_aereo_chart_image(
            gerador_de_graficos_tabelas.gerar_grafico_esforco_aereo_impressao(
                aeronave='E-99',
                grupo='DCTA',
                esforco_aereo_df=esforco_aereo_df)),
        'Esforço Aéreo - R99 COMPREP': generate_esforco_aereo_chart_image(
            gerador_de_graficos_tabelas.gerar_grafico_esforco_aereo_impressao(
                aeronave='R-99',
                grupo='COMPREP',
                esforco_aereo_df=esforco_aereo_df)),
        'Esforço Aéreo - R99 COMAE': generate_esforco_aereo_chart_image(
            gerador_de_graficos_tabelas.gerar_grafico_esforco_aereo_impressao(
                aeronave='R-99',
                grupo='COMAE',
                esforco_aereo_df=esforco_aereo_df)),
        'Esforço Aéreo - R99 DCTA': generate_esforco_aereo_chart_image(
            gerador_de_graficos_tabelas.gerar_grafico_esforco_aereo_impressao(
                aeronave='R-99',
                grupo='DCTA',
                esforco_aereo_df=esforco_aereo_df)),
    }


# Função para criar a apresentação do PowerPoint
def create_presentation(dicionario_de_graficos):
    prs = Presentation("Apresentação semanal - Slides.pptx")
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text in dicionario_de_graficos:
                    image_stream = dicionario_de_graficos[text]
                    slide.shapes.add_picture(image_stream, left=Inches(0.6968), top=Inches(2.177))

    presentation_io = io.BytesIO()
    prs.save(presentation_io)
    presentation_io.seek(0)
    return presentation_io


# Função para criar um link de download
def get_binary_file_downloader_link(binary_file, file_label):
    b64 = base64.b64encode(binary_file.read()).decode('UTF-8')
    href = f'<a href="data:file/pptx;base64,{b64}" download="{file_label}">Clique aqui para baixar a apresentação</a>'
    return href
